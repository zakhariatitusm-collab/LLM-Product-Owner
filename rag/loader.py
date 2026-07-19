from pathlib import Path

from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader

from docx import Document as DocxDocument
import pandas as pd
from pptx import Presentation


def load_pdf(file_path: str):
    loader = PyPDFLoader(file_path)
    return loader.load()


def load_docx(file_path: str):
    doc = DocxDocument(file_path)

    text = "\n".join(
        paragraph.text
        for paragraph in doc.paragraphs
        if paragraph.text.strip()
    )

    return [
        Document(
            page_content=text,
            metadata={"source": file_path}
        )
    ]


def load_txt(file_path: str):
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()

    return [
        Document(
            page_content=text,
            metadata={"source": file_path}
        )
    ]


def load_csv(file_path: str):
    df = pd.read_csv(file_path)

    return [
        Document(
            page_content=df.to_string(index=False),
            metadata={"source": file_path}
        )
    ]


def load_excel(file_path: str):
    df = pd.read_excel(file_path)

    return [
        Document(
            page_content=df.to_string(index=False),
            metadata={"source": file_path}
        )
    ]


def load_pptx(file_path: str):
    prs = Presentation(file_path)

    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return [
        Document(
            page_content="\n".join(text),
            metadata={"source": file_path}
        )
    ]


def load_document(file_path: str):

    extension = Path(file_path).suffix.lower()

    if extension == ".pdf":
        return load_pdf(file_path)

    elif extension == ".docx":
        return load_docx(file_path)

    elif extension == ".txt":
        return load_txt(file_path)

    elif extension == ".csv":
        return load_csv(file_path)

    elif extension == ".xlsx":
        return load_excel(file_path)

    elif extension == ".pptx":
        return load_pptx(file_path)

    else:
        raise ValueError(f"Unsupported file type: {extension}")